from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
import argparse
from pathlib import Path
import zipfile
from xml.etree import ElementTree as ET
import os
import subprocess
import sys


# 如何检查 PPT 是否已有母版
def inspect_master(pptx_path):
    prs = Presentation(pptx_path)
    master = prs.slide_masters[0]
    print(f"母版名称: {master.name}")
    print(f"包含 {len(master.slide_layouts)} 种版式")
    print(f"母版中形状数量: {len(master.shapes)}")  # 包括 Logo、页脚等

    for i, layout in enumerate(master.slide_layouts):
        ph_count = len(layout.placeholders)
        print(f"  版式 {i}: {layout.name} → {ph_count} 个占位符")
        # for ph in layout.placeholders:
        #     # 占位符类型（如 TITLE, BODY, PICTURE 等）
        #     ph_type = ph.placeholder_format.type
        #     ph_type_name = str(ph_type).split('.')[-1] if hasattr(ph_type, 'name') else str(ph_type)
        #
        #     # 占位符索引（用于内部引用）
        #     idx = ph.placeholder_format.idx
        #
        #     # 占位符名称（可自定义，通常为空）
        #     name = ph.name if ph.name else "(unnamed)"
        #
        #     # 位置和大小（EMU 单位，可转换为厘米或英寸）
        #     left, top = ph.left, ph.top
        #     width, height = ph.width, ph.height
        #
        #     print(f"    - 占位符 idx={idx}, 类型={ph_type_name}, 名称='{name}'")
        #     print(f"      位置: ({left}, {top}) | 尺寸: {width}×{height} EMU")
        print("-----------------------------------")


# 中文名 → 英文名 映射表
# 确保模板中包含以下英文布局名（pandoc 依赖它们）：
# Table of Contents pandoc 不会自动生成目录页，也不会查找此 layout
# Vertical Title and Text 等竖排布局 pandoc 从不使用这些特殊布局
# Content 不存在于标准模板中，不要映射 "目录": "Content"
mapping = {
    "标题幻灯片": "Title Slide",
    "标题和内容": "Title and Content",
    "目录": "Title and Content",  # 如果你有“目录”页，通常应使用 "Title and Content"
    "节标题": "Section Header",
    "两栏内容": "Two Content",
    "比较": "Comparison",
    "仅标题": "Title Only",
    "空白": "Blank",
    "内容与标题": "Content with Caption",
    "图片与标题": "Picture with Caption",
    # "标题和竖排文字": "Title and Vertical Text",
    # "竖排标题与文本": "Vertical Title and Text",
}


def generate_english_template(pptx_path, output_path):
    """
    将PPTX文件中的所有版式的名称从中文翻译为英文，并保存为新文件。

    :param pptx_path: 原始PPTX文件路径
    :param output_path: 输出的新PPTX文件路径
    """
    prs = Presentation(pptx_path)

    # 遍历每个母版及其版式
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name in mapping:
                new_name = mapping[layout.name]
                print(f"Renaming layout '{layout.name}' to '{new_name}'")
                layout._element.cSld.attrib['name'] = new_name
            else:
                print(f"can not mapping '{layout.name}'")

    # 保存新的PPT文件
    prs.save(output_path)
    print(f"New template saved to {output_path}")


if __name__ == "__main__":
    pptx_path = "/Users/jiaxiaopeng/ppt/source.pptx"
    inspect_master(pptx_path)
    out_path = "/Users/jiaxiaopeng/ppt/my_t.pptx"
    generate_english_template(pptx_path, out_path);
